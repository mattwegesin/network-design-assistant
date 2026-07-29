import os
import base64
from flask import Flask, render_template, request, jsonify, flash
from werkzeug.utils import secure_filename
from google import genai
from google.genai import types
from dotenv import load_dotenv
import markdown
import pandas as pd
import docx2txt
from pptx import Presentation

load_dotenv()

app = Flask(__name__)
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'super-secret-key-123')

# Check for API Key
if not os.environ.get("GEMINI_API_KEY"):
    print("WARNING: GEMINI_API_KEY environment variable not set.")

# --- System Instructions ---
NETWORK_DESIGNER_INSTRUCTION = """# Network Designer (Master Engine)

This skill acts as a Network Design Engineer. It generates professional Wi-Fi network designs, calculates switch port counts, determines hardware requirements, and produces a Bill of Materials (BOM) by applying standardized design logic to brand-specific parameters.

## Workflow

### 1. Input Analysis (AP Layout)
If the user provides an **AP Layout** (image, PDF, or description), prioritize the visual placement over standard scaling rules. Analyze the layout to count APs based on the following color-coding standards:
- **Blue APs**: Common Area / Hallway (Models: Hallway/Common Area APs from reference)
- **Purple APs**: In-room (Models: In-room APs from reference)
- **Orange APs**: Outdoor (Models: Outdoor APs from reference)
- **Green APs**: Meeting Area (Models: High-density/Meeting Area APs from reference)

### 2. Standardized Design Logic
Apply the following rules using the values found in the provided reference context:

#### A. Access Point Deployment
- **Guest Rooms:**
  - If walls are **Sheetrock/Drywall**: Deploy 1 Guest Room AP per 2 rooms (round up).
  - If walls are **Concrete/Brick**: Deploy 1 Guest Room AP per 1 room.
  - *Exception:* If the brand specifies a "Hallway Design" and room drops are missing, follow the specific Hallway scaling rules in the reference file.
- **Common Areas:** Deploy 1 AP per specified area (Lobby, Fitness, Meeting, etc.) using the models defined in the reference file.
- **Mounting:** Include the specific mounting hardware SKU for every AP as defined in the reference file.

#### B. Switch Port Calculation (Per Equipment Room)
1. **Base Ports:** Count 1 port per Guest Room drop + 1 port per Common Area AP.
2. **Buffer:** Apply a **15% buffer** (Total * 1.15).
3. **Switch Selection:** 
   - Primary: Use the **48-Port Switch** model.
   - Optimization: If using a **24-Port Switch** model would result in fewer than 24 wasted ports compared to a 48-port unit, use the 24-port model.
   - Goal: Minimize wasted ports while ensuring 15% growth capacity.

#### C. PoE Budget Validation
1. Calculate the total PoE load for all powered devices (APs, SFPs) in the room using the **PoE Load** values in the reference file.
2. Compare the total against **85% of the Switch's Max PoE Budget**.
3. If the load exceeds the 85% threshold, add an additional switch to split the load.

#### D. Connectivity & Infrastructure
- **Cabling:** Allocate **200 ft** of Indoor/Outdoor CAT6 cable per AP as defined. Include 6-inch patch cables for in-room APs and IPTV/Phones.
- **Fiber:** Include 2 SFPs per fiber run. Match the SFP model to the fiber type (Multimode/Singlemode) defined in the reference.
- **Racks/UPS:** Include 1 x 12U Rack for MDF and 1 x 8U Rack for each IDF. Include 1 x 1500VA UPS per rack.

#### E. Labor & Fees
- **Calculated Labor:** Use the labor SKUs from the reference file.
- **Counters:** 
  - Install Labor: Count per AP, per Switch, and per UPS.
  - Project Management/Documentation: Apply standard brand SKUs.
- **Monthly Fees:** Include the required Support/MSF SKUs per room or per property.

### 3. Generate Output
Produce a comprehensive Network Design Summary and a structured Bill of Materials (BOM) table.
- **Show your work:** Provide a step-by-step breakdown of calculations (APs per floor, Ports per IDF, PoE Load validation).
- **Tabular BOM:** Present the final hardware list as a Markdown table containing Part Number, Description, and Quantity.

## Design Constraints
- **Data-Driven:** Always fetch values (Part Numbers, PoE loads, SKUs) from the provided reference context.
- **Modular:** Perform all switch and PoE calculations on a per-equipment-room (MDF/IDF) basis.
- **Precision:** Never round until the final buffer calculation. Always round UP for hardware counts.
- **Conservative Approach:** Always prioritize higher AP counts over lower, and always use per-floor rounding (rounding up) for hallway density calculations to ensure 100% coverage.
"""

READINESS_INSTRUCTION = """# Network Design Readiness Analyst

This skill enables you to act as a technical intake specialist for Wi-Fi engineering projects. You will audit provided files/images to ensure all variables required for RF design and equipment selection are present.

## Workflow

### 1. Initial Audit
1.  **Inventory Files**: Analyze all provided images/documents.
2.  **Classify Documents**: Identify Floorplans, Photos, and Notes.
3.  **Variable Extraction**: Attempt to extract the following from document contents:
    - Room counts per floor.
    - Wall materials (Concrete, Drywall, etc.).
    - Status of existing wired drops.
    - Common area requirements.
    - **AP Layout Presence**: Check if a color-coded AP layout is present (Blue: Hallway, Purple: In-Room, Orange: Outdoor, Green: Meeting).

### 2. Readiness Reporting
Provide a structured report using the following format:

**Project Name:** [Derived from context]
**Overall Status:** [READY / MISSING INFORMATION]

**Data Point Checklist:**
- [ ] Floorplans (Scalable)
- [ ] Wall Materials Identified
- [ ] Room Count Verified
- [ ] Wired Drop Status
- [ ] MDF/IDF Locations
- [ ] Common Area Scope
- [ ] **Color-Coded AP Layout** (If present, design can proceed immediately)

**Findings:**
- List what was found.
- Identify specific missing variables.
- Note if a color-coded AP layout was found.

**Recommendations:**
- If READY: Propose starting the design/BOM generation.
- If MISSING INFO: Draft a concise request for the user to provide the specific missing data points.
"""

def get_reference_content(brand_file):
    filepath = os.path.join('references', brand_file)
    try:
        with open(filepath, 'r') as f:
            return f.read()
    except FileNotFoundError:
        return ""

def process_file(file):
    filename = file.filename.lower()
    
    # Process Excel files to CSV text
    if filename.endswith('.xlsx') or filename.endswith('.xls'):
        try:
            df = pd.read_excel(file)
            csv_data = df.to_csv(index=False)
            return {
                "mime_type": "text/csv",
                "data": csv_data.encode('utf-8')
            }
        except Exception as e:
            raise ValueError(f"Failed to extract data from Excel file ({filename}): {str(e)}")
            
    # Process Word documents to plain text
    elif filename.endswith('.docx'):
        try:
            text = docx2txt.process(file)
            return {
                "mime_type": "text/plain",
                "data": text.encode('utf-8')
            }
        except Exception as e:
            raise ValueError(f"Failed to extract text from Word document ({filename}): {str(e)}")
            
    # Process PowerPoint presentations to plain text
    elif filename.endswith('.pptx'):
        try:
            prs = Presentation(file)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return {
                "mime_type": "text/plain",
                "data": ("\n".join(text_runs)).encode('utf-8')
            }
        except Exception as e:
            raise ValueError(f"Failed to extract text from PowerPoint ({filename}): {str(e)}")

    # Fallback for PDFs, TXT, CSV, and Images (or if parsing fails)
    file.seek(0)
    bytes_data = file.read()
    file.seek(0)
    
    mime_type = file.mimetype
    # Fix generic octet-stream mimetypes for known text/csv files
    if mime_type == 'application/octet-stream' or not mime_type:
        if filename.endswith('.csv'):
            mime_type = 'text/csv'
        elif filename.endswith('.txt') or filename.endswith('.md'):
            mime_type = 'text/plain'
            
    return {
        "mime_type": mime_type,
        "data": bytes_data
    }

@app.route('/')
def index():
    brands = [
        {"filename": "best-western-aruba.md", "name": "Best Western Aruba"},
        {"filename": "best-western-omada.md", "name": "Best Western Omada"},
        {"filename": "choice-aruba.md", "name": "Choice Aruba"},
        {"filename": "ihg-meraki.md", "name": "IHG Meraki"},
        {"filename": "independent-omada.md", "name": "Independent Omada"},
        {"filename": "wyndham-aruba.md", "name": "Wyndham Aruba"},
        {"filename": "design-criteria.md", "name": "General Design Criteria"}
    ]
    return render_template('index.html', brands=brands)

@app.route('/analyze', methods=['POST'])
def analyze():
    try:
        if 'GEMINI_API_KEY' not in os.environ:
            return jsonify({'error': 'API Key not configured on the server.'}), 500

        task_type = request.form.get('taskType')
        brand_file = request.form.get('brandFile')
        include_sow = request.form.get('includeSow') == 'true'
        include_diagram = request.form.get('includeDiagram') == 'true'
        prompt = request.form.get('prompt', '')

        if not task_type:
            return jsonify({'error': 'Task type is required.'}), 400

        # Build the system instruction and context
        system_instruction = ""
        context = ""

        if task_type == 'designer':
            system_instruction = NETWORK_DESIGNER_INSTRUCTION
            if include_sow:
                system_instruction += """

        ### Statement of Work (SOW) Generation
        You must also generate a Statement of Work based on the standard format. 
        Maintain the professional formatting:
        - Executive Summary
        - Scope of Work
        - Deliverables
        - Timeline
        - Exclusions
        - Pricing Structure

        Crucially, populate **Section 3.1 (Equipment)** with the specific hardware required for each equipment room (MDF/IDF) based on your network design BOM.
        Format for Section 3.1:
        #### 3.1 Equipment per Location
        - **MDF (Main Distribution Frame)**
        - [Quantity] x [Part Number] - [Description]
        - **IDF [Number] (Intermediate Distribution Frame)**
        - [Quantity] x [Part Number] - [Description]
        """
            if include_diagram:
                system_instruction += """

        ### Logical Network Diagram
        You must generate a logical network diagram using Mermaid.js syntax. 
        Use a Markdown code block with `mermaid` as the language (i.e. ```mermaid).
        Map the connections from the Gateway -> MDF -> IDFs -> Switches -> APs.
        Keep it high-level and clear.
        """

            if brand_file:
                context = f"\n\n--- REFERENCE DATA ({brand_file}) ---\n" + get_reference_content(brand_file)
        elif task_type == 'readiness':
            system_instruction = READINESS_INSTRUCTION
            # For readiness, we load general design criteria if present, or just pass context
            context = f"\n\n--- REFERENCE DATA (design-criteria.md) ---\n" + get_reference_content('design-criteria.md')

        # Prepare contents for Gemini
        contents = []
        
        # Process uploaded files (images/pdfs/office)
        uploaded_files = request.files.getlist('files')
        
        for file in uploaded_files:
            if file.filename:
                file_part = process_file(file)
                contents.append(types.Part.from_bytes(
                    data=file_part['data'],
                    mime_type=file_part['mime_type'],
                ))

        # Build the final text prompt
        final_prompt = prompt
        if context:
             final_prompt += context
             
        contents.append(final_prompt)

        # Call Gemini API
        api_key = os.environ.get('GEMINI_API_KEY')
        if not api_key:
            return jsonify({'error': 'API Key missing on server.'}), 500
            
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(
            model='gemini-3.1-pro-preview',
            contents=contents,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                temperature=0.2, # Low temperature for more deterministic/factual output
            )
        )
        
        # Convert Markdown response to HTML for display
        html_content = markdown.markdown(response.text, extensions=['tables'])
        
        return jsonify({'result': html_content, 'raw_markdown': response.text})
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=True)
